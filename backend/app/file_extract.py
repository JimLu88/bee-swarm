"""v6-Y 文档解析: 把用户上传的 base64 文件 (xlsx/pdf/docx/pptx/csv/txt...) 转成纯文字.

设计要点:
- 进程内解析 (不依赖 bee-light-exec 跨服务), 转文字后拼进 task.
- 文字对所有模型免费可读 (含瞎子 deepseek/ollama), 不走视觉, 也不用按档区分.
- 单文件 + 总量都有上限, 防止 token 爆炸; 超限截断并明确提示.
"""
from __future__ import annotations

import base64
import csv
import io
import json
from typing import Any

# 单文件最多保留多少字 (解析后文字); 总量上限在 extract_files 里
_PER_FILE_MAX_CHARS = 8000
_TOTAL_MAX_CHARS = 12000
_XLSX_MAX_ROWS = 200
_XLSX_MAX_COLS = 30
_PDF_MAX_PAGES = 30
_PPTX_MAX_SLIDES = 60


def _truncate(text: str, limit: int) -> str:
    if len(text) <= limit:
        return text
    return text[:limit] + f"\n…[已截断, 原文 {len(text)} 字, 只取前 {limit} 字]"


def _ext_of(name: str) -> str:
    name = (name or "").lower().strip()
    if "." not in name:
        return ""
    return name.rsplit(".", 1)[-1]


def _decode_b64(content_b64: str) -> bytes:
    s = content_b64 or ""
    # 容忍 data URL 前缀 (data:application/...;base64,XXXX)
    if s.startswith("data:") and "," in s:
        s = s.split(",", 1)[1]
    return base64.b64decode(s)


# ---- 各类型解析器: 输入 bytes, 输出文字 ----

def _parse_xlsx(raw: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    out: list[str] = []
    for ws in wb.worksheets:
        out.append(f"## 工作表: {ws.title}")
        rows_done = 0
        for row in ws.iter_rows(values_only=True):
            if rows_done >= _XLSX_MAX_ROWS:
                out.append(f"…[超过 {_XLSX_MAX_ROWS} 行, 已截断]")
                break
            cells = ["" if c is None else str(c) for c in row[:_XLSX_MAX_COLS]]
            if any(cells):
                out.append(" | ".join(cells))
                rows_done += 1
    wb.close()
    return "\n".join(out)


def _parse_csv(raw: bytes) -> str:
    text = raw.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    out: list[str] = []
    for i, row in enumerate(reader):
        if i >= _XLSX_MAX_ROWS:
            out.append(f"…[超过 {_XLSX_MAX_ROWS} 行, 已截断]")
            break
        out.append(" | ".join(row[:_XLSX_MAX_COLS]))
    return "\n".join(out)


def _parse_pdf(raw: bytes) -> str:
    from pypdf import PdfReader
    rdr = PdfReader(io.BytesIO(raw))
    out: list[str] = []
    for i, page in enumerate(rdr.pages):
        if i >= _PDF_MAX_PAGES:
            out.append(f"…[超过 {_PDF_MAX_PAGES} 页, 已截断]")
            break
        txt = (page.extract_text() or "").strip()
        if txt:
            out.append(f"[第{i+1}页]\n{txt}")
    result = "\n\n".join(out)
    if not result.strip():
        return "[PDF 没解析出文字 — 可能是扫描件/图片型 PDF, 建议改用图片上传走视觉识别]"
    return result


def _parse_docx(raw: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(raw))
    out: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            out.append(p.text)
    # 表格也抓
    for t in doc.tables:
        for row in t.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                out.append(" | ".join(cells))
    return "\n".join(out)


def _parse_pptx(raw: bytes) -> str:
    from pptx import Presentation
    pres = Presentation(io.BytesIO(raw))
    out: list[str] = []
    for i, slide in enumerate(pres.slides):
        if i >= _PPTX_MAX_SLIDES:
            out.append(f"…[超过 {_PPTX_MAX_SLIDES} 页, 已截断]")
            break
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
        if parts:
            out.append(f"[幻灯片{i+1}]\n" + "\n".join(parts))
    return "\n\n".join(out)


def _parse_text(raw: bytes) -> str:
    return raw.decode("utf-8", errors="replace")


def _parse_json(raw: bytes) -> str:
    try:
        obj = json.loads(raw.decode("utf-8", errors="replace"))
        return json.dumps(obj, ensure_ascii=False, indent=2)
    except Exception:
        return raw.decode("utf-8", errors="replace")


# 扩展名 → 解析器
_PARSERS = {
    "xlsx": _parse_xlsx, "xlsm": _parse_xlsx,
    "csv": _parse_csv, "tsv": _parse_csv,
    "pdf": _parse_pdf,
    "docx": _parse_docx,
    "pptx": _parse_pptx,
    "txt": _parse_text, "md": _parse_text, "log": _parse_text,
    "json": _parse_json,
}

# 这些扩展名应当走图片视觉路径, 不在这里解析
IMAGE_EXTS = {"png", "jpg", "jpeg", "gif", "webp", "bmp"}


def extract_one(name: str, content_b64: str) -> str:
    """解析单个文件 → 文字. 失败返回 [解析失败:...] 而非抛异常 (一个坏文件不该挂掉整次决策)."""
    ext = _ext_of(name)
    if ext in IMAGE_EXTS:
        return f"[{name} 是图片 — 应走图片上传走视觉, 这里跳过]"
    parser = _PARSERS.get(ext)
    if not parser:
        # 未知类型: 尝试当文本读
        try:
            return _truncate(_parse_text(_decode_b64(content_b64)), _PER_FILE_MAX_CHARS)
        except Exception as e:
            return f"[不支持的文件类型 .{ext}: {name} ({e!r})]"
    try:
        raw = _decode_b64(content_b64)
        return _truncate(parser(raw), _PER_FILE_MAX_CHARS)
    except Exception as e:
        return f"[解析失败 {name}: {e!r}]"


def extract_files(files: list[dict[str, Any]] | None) -> str:
    """解析一组文件 → 拼成一个 [附件内容] 文字块 (拼进 task 给所有部门看).

    files: [{name, content_b64}, ...]
    返回: 空串 (没文件时) 或 "[用户上传的附件]\n### 文件名\n...内容..."
    """
    if not files:
        return ""
    blocks: list[str] = []
    used = 0
    for f in files:
        name = str(f.get("name") or "未命名文件")
        b64 = str(f.get("content_b64") or "")
        if not b64:
            continue
        text = extract_one(name, b64)
        remaining = _TOTAL_MAX_CHARS - used
        if remaining <= 0:
            blocks.append(f"### {name}\n[附件总量超限, 此文件及之后未解析]")
            break
        if len(text) > remaining:
            text = _truncate(text, remaining)
        used += len(text)
        blocks.append(f"### {name}\n{text}")
    if not blocks:
        return ""
    return "[用户上传的附件 — 请结合分析]\n" + "\n\n".join(blocks)
